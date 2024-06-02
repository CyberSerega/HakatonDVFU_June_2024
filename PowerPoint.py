import os
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from ImageColor import generate_color
from tkinter import Tk
from tkinter.filedialog import asksaveasfilename
from TextGenerate import make_text

def generate_slide1(prs, layout, title, subtitle,  image_path, url, color_combinations, font):
    # Создание слайда
    slide = prs.slides.add_slide(layout)
    # Назначение цвета фона для всего слайда
    slide_background = color_combinations[0]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*slide_background)

    # Добавление заголовка
    title_box = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(15.24), Cm(1.93))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.font.bold = True
    p.font.size = Pt(39)
    p.font.name = font
    p.text = title

    # Добавление подзаголовка
    #subtitle_box = slide.shapes.add_textbox(Cm(1), Cm(7), Cm(15.24), Cm(1.93))
    #subtitle_frame = subtitle_box.text_frame
    #p = subtitle_frame.add_paragraph()
    #p.font.size = Pt(23)
    #p.text = subtitle
    #subtitle_frame.paragraphs[0].font.name = font

    # Добавление заголовка в закругленную форму
    #title_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3), Cm(15.24), Cm(1.93))
    #title_shape.text = title
    #title_frame = title_shape.text_frame
    #title_frame.word_wrap = True
    #p = title_frame.paragraphs[0]  # Получаем первый параграф
    #p.font.bold = True
    #p.font.size = Pt(39)
    #p.font.name = font

    # Добавление подзаголовка в закругленную форму
    subtitle_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(7.5), Cm(15.24), Cm(1.93))
    subtitle_shape.text = subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.paragraphs[0].font.name = font
    subtitle_frame.paragraphs[0].font.size = Pt(23)

    # Добавление объекта для изображения
    slide.shapes.add_picture(image_path, Cm(17.5), Cm(0), Cm(14.25), Cm(15.69))

    # Добавление кнопки с гиперссылкой
    if url == "./res/qr-code.png":
        slide.shapes.add_picture("./res/qr-code.png", Cm(1), Cm(10), width=Cm(5))
    elif url is not None:
        button_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(12), Cm(8.81), Cm(2.3))
        button_shape.text = "Узнать подробнее"
        link = button_shape.click_action.hyperlink
        link.address = url
        text_frame = button_shape.text_frame  # Получение объекта текстового кадра кнопки
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Установка выравнивания текста
        paragraph = text_frame.paragraphs[0]  # Получение объекта параграфа кнопки
        paragraph.font.size = Pt(24)  # Установка размера шрифта

        button_color, button_text_color = color_combinations[4], color_combinations[0]
        button_shape.fill.solid()
        button_shape.fill.fore_color.rgb = RGBColor(*button_color)
        button_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*button_text_color)

    # Назначение цвета фона и текста для заголовка и подзаголовка
    print(color_combinations)
    title_color, title_text_color = color_combinations[0], color_combinations[1]
    subtitle_color, subtitle_text_color = color_combinations[2], color_combinations[3]

    # Установка цвета фона и текста
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(*title_color)
    p.font.color.rgb = RGBColor(*title_text_color)  # Применение цвета к тексту


    #subtitle_box.fill.solid()
    #subtitle_box.fill.fore_color.rgb = RGBColor(*subtitle_color)
    #subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*subtitle_text_color)

    #title_shape.fill.solid()
    #title_shape.fill.fore_color.rgb = RGBColor(*title_color)
    #title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_text_color)

    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = RGBColor(*subtitle_color)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*subtitle_text_color)

def generate_slide2(prs, layout, title, subtitle,overtitle, image_path, url, color_combinations, font):
    # Создание слайда
    slide = prs.slides.add_slide(layout)
    # Назначение цвета фона для всего слайда
    slide_background = color_combinations[0]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*slide_background)

    # Надзаголовок
    overtitle_box = slide.shapes.add_textbox(Cm(1), Cm(2), Cm(15.24), Cm(1.93))
    overtitle_frame = overtitle_box.text_frame
    p1 = overtitle_frame.add_paragraph()
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.text = overtitle

    # Добавление заголовка
    title_box = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(16.5), Cm(4.45))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p2 = title_frame.add_paragraph()
    p2.font.bold = True
    p2.font.size = Pt(40)
    p2.font.name = font
    p2.text = title


    # Добавление подзаголовка
    #subtitle_box = slide.shapes.add_textbox(Cm(1), Cm(7), Cm(15.24), Cm(1.93))
    #subtitle_frame = subtitle_box.text_frame
    #subtitle_frame.word_wrap = True
    #p = subtitle_frame.add_paragraph()
    #p.font.size = Pt(24)
    #p.text = subtitle
    #subtitle_frame.paragraphs[0].font.name = font

    # Добавление подзаголовка в закругленную форму
    subtitle_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(7.5), Cm(15.24), Cm(1.93))
    subtitle_shape.text = subtitle
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.paragraphs[0].font.name = font
    subtitle_frame.paragraphs[0].font.size = Pt(24)

    # Добавление объекта для изображения
    slide.shapes.add_picture(image_path, Cm(17.5), Cm(0), Cm(14.25), Cm(15.69))

    # Добавление кнопки с гиперссылкой
    if url == "./res/qr-code.png":
        slide.shapes.add_picture("./res/qr-code.png", Cm(1), Cm(10), width=Cm(5))
    elif url is not None:
        button_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(12), Cm(6.35), Cm(2.3))
        button_shape.text = "Подробнее"
        link = button_shape.click_action.hyperlink
        link.address = url
        text_frame = button_shape.text_frame # Получение объекта текстового кадра кнопки
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER # Установка выравнивания текста
        paragraph = text_frame.paragraphs[0] # Получение объекта параграфа кнопки
        paragraph.font.size = Pt(24)  # Установка размера шрифта

        button_color, button_text_color = color_combinations[4], color_combinations[0]
        button_shape.fill.solid()
        button_shape.fill.fore_color.rgb = RGBColor(*button_color)
        button_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*button_text_color)

    print(color_combinations)
    # Назначение цвета фона и текста для заголовка и подзаголовка
    title_color, title_text_color = color_combinations[0], color_combinations[1]
    subtitle_color, subtitle_text_color = color_combinations[2], color_combinations[3]

    #title_box.fill.solid()
    #title_box.fill.fore_color.rgb = RGBColor(*title_color)
    p2.font.color.rgb = RGBColor(*title_text_color)
    p1.font.color.rgb = RGBColor(*title_text_color)


    #title_box.fill.solid()
    #title_box.fill.fore_color.rgb = RGBColor(*title_color)
    #title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_text_color)

    #subtitle_box.fill.solid()
    #subtitle_box.fill.fore_color.rgb = RGBColor(*subtitle_color)
    #subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*subtitle_text_color)

    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = RGBColor(*subtitle_color)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*subtitle_text_color)


def generate_slide3(prs, layout, title, subtitle, image_path, url, color_combinations, font):
    # Создание слайда
    slide = prs.slides.add_slide(layout)
    # Назначение цвета фона для всего слайда
    slide_background = color_combinations[0]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*slide_background)

    # Добавление заголовка
    title_box = slide.shapes.add_textbox(Cm(1), Cm(4), Cm(15.24), Cm(1.93))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.add_paragraph()
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.name = font
    p.text = title

    # Добавление объекта для изображения
    slide.shapes.add_picture(image_path, Cm(17.5), Cm(0), Cm(14.25), Cm(15.69))

    # Добавление кнопки с гиперссылкой
    if url == "./res/qr-code.png":
        slide.shapes.add_picture("./res/qr-code.png", Cm(1), Cm(10), width=Cm(5))
    elif url is not None:
        button_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(11), Cm(8.81), Cm(2.3))
        button_shape.text = "Узнать подробнее"
        link = button_shape.click_action.hyperlink
        link.address = url
        text_frame = button_shape.text_frame  # Получение объекта текстового кадра кнопки
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Установка выравнивания текста
        paragraph = text_frame.paragraphs[0]  # Получение объекта параграфа кнопки
        paragraph.font.size = Pt(24)  # Установка размера шрифта

        button_color, button_text_color = color_combinations[4], color_combinations[0]
        button_shape.fill.solid()
        button_shape.fill.fore_color.rgb = RGBColor(*button_color)
        button_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*button_text_color)

    # Назначение цвета фона и текста для заголовка и подзаголовка
    title_color, title_text_color = color_combinations[0], color_combinations[1]

    p.font.color.rgb = RGBColor(*title_text_color)

    #title_box.fill.solid()
    #title_box.fill.fore_color.rgb = RGBColor(*title_color)
    #title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_text_color)

def generate_presentation(main_entry, font, switch_1, url):
    prs = Presentation()

    prs.slide_width = Inches(1200 / 96)
    prs.slide_height = Inches(593 / 96)


    
    # 1 слайд
    slide_layout = prs.slide_layouts[6]
    array = make_text(main_entry, 2)
    color_combinations = generate_color("./res/1.jpg")
    generate_slide1(prs, slide_layout, array[0], array[1], "./res/1.jpg", url,
                   color_combinations, font)
    # 2 слайд
    slide_layout = prs.slide_layouts[6]
    array = make_text(main_entry, 3)
    color_combinations = generate_color("./res/2.jpg")
    generate_slide2(prs, slide_layout, array[0], array[1], array[2], "./res/2.jpg", url,
                   color_combinations, font)

    # 3 слайд
    slide_layout = prs.slide_layouts[6]
    array = make_text(main_entry, 1)
    color_combinations = generate_color("./res/3.jpg")
    generate_slide3(prs, slide_layout, array, "Подзаголовок", "./res/3.jpg", url,
                    color_combinations, font)

    # Используем tkinter для запроса пути сохранения файла
    root = Tk()
    root.withdraw()  # Скрываем основное окно
    file_path = asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])

    if file_path:
        prs.save(file_path)
        print(f"Презентация сохранена по пути: {file_path}")
        os.startfile(file_path)
    else:
        print("Сохранение отменено пользователем")
